#!/usr/bin/env python3
"""Leitor universal de arquivos — extrai texto de PDF, DOCX, DOC, ODT, ODS, ODP,
PPTX, XLSX, imagens (OCR), CSV, JSON, HTML, RTF e mais."""
import sys, os, json, csv, io, re, mimetypes
from pathlib import Path
from typing import Optional

SUPPORTED_EXTENSIONS = {
    '.pdf', '.docx', '.doc', '.odt', '.ods', '.odp', '.rtf',
    '.pptx', '.xlsx', '.xls', '.csv', '.html', '.htm',
    '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.tif', '.webp',
}

TEXT_EXTS = {
    '.md', '.txt', '.py', '.js', '.ts', '.json', '.yaml', '.yml',
    '.toml', '.cfg', '.ini', '.sql', '.xml', '.log', '.env',
    '.sh', '.bash', '.zsh', '.css', '.scss', '.less', '.vue',
    '.svelte', '.jsx', '.tsx', '.rb', '.go', '.rs', '.java',
    '.c', '.cpp', '.h', '.swift', '.kt', '.dart', '.php', '.r', '.lua',
    '.svg',
}

# ============================================================
# PDF
# ============================================================

def extract_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
        r = PdfReader(str(path))
        pages = []
        for p in r.pages:
            t = p.extract_text()
            if t and t.strip():
                pages.append(t.strip())
        if pages:
            return '\n\n'.join(pages)
    except Exception:
        pass

    # Fallback: pdftotext
    try:
        import subprocess
        r = subprocess.run(
            ['pdftotext', '-layout', str(path), '-'],
            capture_output=True, text=True, timeout=30,
        )
        if r.stdout and r.stdout.strip():
            return r.stdout
    except Exception:
        pass

    # Fallback: OCR (scanned PDF)
    try:
        import subprocess
        from PIL import Image
        import pytesseract
        r = subprocess.run(
            ['pdftoppm', '-png', '-r', '300', str(path)],
            capture_output=True, timeout=60,
        )
        if r.stdout:
            img = Image.open(io.BytesIO(r.stdout))
            return pytesseract.image_to_string(img, lang='por+eng')
    except Exception:
        pass

    return '[PDF: nao foi possivel extrair texto]'

# ============================================================
# Word (.docx e .doc)
# ============================================================

def extract_docx(path: Path) -> str:
    from docx import Document
    d = Document(str(path))
    lines = []
    for p in d.paragraphs:
        t = p.text.strip()
        if t:
            lines.append(t)
    for t in d.tables:
        for row in t.rows:
            cells = [c.text.strip() for c in row.cells]
            lines.append(' | '.join(cells))
    return '\n'.join(lines)

def extract_doc(path: Path) -> str:
    """Word antigo (.doc) via antiword ou olefile + zlib."""
    try:
        import subprocess
        r = subprocess.run(
            ['antiword', str(path)],
            capture_output=True, text=True, timeout=15,
        )
        if r.stdout and r.stdout.strip():
            return r.stdout
    except Exception:
        pass

    try:
        import subprocess
        r = subprocess.run(
            ['catdoc', str(path)],
            capture_output=True, text=True, timeout=15,
        )
        if r.stdout and r.stdout.strip():
            return r.stdout
    except Exception:
        pass

    # Fallback: olefile + zlib (texto cru)
    try:
        import olefile, zlib
        ole = olefile.OleFileIO(str(path))
        if ole.exists('WordDocument'):
            stream = ole.openstream('WordDocument')
            data = stream.read()
            # Tenta extrair texto Unicode
            if ole.exists('1Table') or ole.exists('0Table'):
                text = data.decode('utf-16-le', errors='replace')
                # Filtra caracteres nao-imprimiveis
                text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)
                if text.strip():
                    return text.strip()
    except Exception:
        pass

    return '[DOC: nao foi possivel extrair texto]'

# ============================================================
# LibreOffice (.odt, .ods, .odp)
# ============================================================

def extract_odt(path: Path) -> str:
    from odf.opendocument import load
    from odf.text import P
    doc = load(str(path))
    texts = []
    for p in doc.getElementsByType(P):
        t = ''.join(node.data for node in p.childNodes if hasattr(node, 'data'))
        if t.strip():
            texts.append(t.strip())
    return '\n'.join(texts)

def extract_ods(path: Path) -> str:
    from odf.opendocument import load
    from odf.table import Table, TableRow, TableCell
    from odf.text import P
    doc = load(str(path))
    lines = []
    for table in doc.getElementsByType(Table):
        name = table.getAttribute('table:name') or ''
        lines.append(f'--- {name} ---')
        for row in table.getElementsByType(TableRow):
            cells = []
            for cell in row.getElementsByType(TableCell):
                texts = []
                for p in cell.getElementsByType(P):
                    t = ''.join(node.data for node in p.childNodes if hasattr(node, 'data'))
                    texts.append(t.strip())
                cells.append(' '.join(texts))
            lines.append('\t'.join(cells))
    return '\n'.join(lines)

def extract_odp(path: Path) -> str:
    from odf.opendocument import load
    from odf.text import P
    doc = load(str(path))
    texts = []
    for p in doc.getElementsByType(P):
        t = ''.join(node.data for node in p.childNodes if hasattr(node, 'data'))
        if t.strip():
            texts.append(t.strip())
    return '\n'.join(texts)

# ============================================================
# RTF
# ============================================================

def extract_rtf(path: Path) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
        text = rtf_to_text(path.read_text(encoding='utf-8', errors='replace'))
        if text.strip():
            return text
    except Exception:
        pass
    return '[RTF: nao foi possivel extrair texto]'

# ============================================================
# PowerPoint
# ============================================================

def extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        lines.append(t)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    lines.append(' | '.join(cells))
    return '\n'.join(lines)

# ============================================================
# Excel
# ============================================================

def extract_xlsx(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), data_only=True)
    lines = []
    for sheet in wb.worksheets:
        lines.append(f'--- {sheet.title} ---')
        for row in sheet.iter_rows(values_only=True):
            lines.append('\t'.join(str(c) if c is not None else '' for c in row))
    return '\n'.join(lines)

def extract_xls(path: Path) -> str:
    try:
        import subprocess
        r = subprocess.run(
            ['ssconvert', '--export-type=Gnumeric_stf:stf_csv', str(path), 'fd://1'],
            capture_output=True, text=True, timeout=30,
        )
        if r.stdout and r.stdout.strip():
            return r.stdout
    except Exception:
        pass
    return '[XLS: nao foi possivel extrair texto — converta para .xlsx]'

# ============================================================
# Imagem / Visao (modelo de visao gratuito) + OCR fallback
# ============================================================

def extract_image_vision(path: Path) -> str:
    """Usa um modelo de visao (via LiteLLM/OpenAI-compatible) para descrever
    a imagem em texto. Requer VISION_MODEL registrado e OPENAI_BASE_URL/OPENAI_API_KEY."""
    try:
        import base64, json, mimetypes, urllib.request
        base = (os.environ.get('OPENAI_BASE_URL') or '').rstrip('/')
        if base and not base.endswith('/v1'):
            base += '/v1'
        if not base:
            return ''
        url = base + '/chat/completions'
        key = os.environ.get('OPENAI_API_KEY', '')
        model = os.environ.get('VISION_MODEL', 'vision-free')
        mime = mimetypes.guess_type(str(path))[0] or 'image/png'
        b64 = base64.b64encode(path.read_bytes()).decode('ascii')
        data_url = f'data:{mime};base64,{b64}'
        prompt = ('Descreva este arquivo/imagem em detalhes para que um assistente de texto '
                  'possa entender seu conteudo completo: transcreva todo o texto visivel, '
                  'tabelas, numeros, elementos visuais e o contexto. Se for um print, diagraMa '
                  'ou captura, seja o mais fiel e completo possivel.')
        payload = {
            'model': model,
            'messages': [{'role': 'user', 'content': [
                {'type': 'text', 'text': prompt},
                {'type': 'image_url', 'image_url': {'url': data_url}},
            ]}],
            'max_tokens': 4000,
        }
        req = urllib.request.Request(
            url, data=json.dumps(payload).encode('utf-8'),
            headers={'Authorization': f'Bearer {key}', 'Content-Type': 'application/json'},
        )
        with urllib.request.urlopen(req, timeout=90) as resp:
            j = json.load(resp)
        return j['choices'][0]['message']['content'].strip()
    except Exception:
        return ''

def extract_image(path: Path) -> str:
    # 1) Tenta visao (modelo gratuito) — leitura semanticamente rica
    vision = extract_image_vision(path)
    if vision:
        return vision
    # 2) Fallback: OCR (se tesseract instalado)
    try:
        import pytesseract
        from PIL import Image, ImageEnhance, ImageFilter
        img = Image.open(str(path))
        img = img.convert('L')
        img = ImageEnhance.Contrast(img).enhance(2.0)
        img = ImageEnhance.Sharpness(img).enhance(2.0)
        text = pytesseract.image_to_string(img, lang='por+eng')
        if text.strip():
            return text
        text = pytesseract.image_to_string(img, lang='por+eng', config='--psm 6 --oem 3')
        return text if text.strip() else '[OCR: nenhum texto detectado]'
    except ImportError:
        return '[Visao e OCR indisponiveis — adicione um modelo de visao gratuito (VISION_MODEL) no LiteLLM]'
    except Exception as e:
        return f'[Erro ao ler imagem: {e}]'

# ============================================================
# CSV
# ============================================================

def extract_csv(path: Path) -> str:
    with open(str(path), newline='', encoding='utf-8', errors='replace') as f:
        reader = csv.reader(f)
        return '\n'.join('\t'.join(row) for row in reader)

# ============================================================
# HTML
# ============================================================

def extract_html(path: Path) -> str:
    from html.parser import HTMLParser
    class T(HTMLParser):
        def __init__(self):
            super().__init__()
            self.text = []
        def handle_data(self, data):
            t = data.strip()
            if t:
                self.text.append(t)
    p = T()
    p.feed(path.read_text(encoding='utf-8', errors='replace'))
    return '\n'.join(p.text)

# ============================================================
# Mapeamento de extratores
# ============================================================

EXTRACTORS = {
    '.pdf': extract_pdf,
    '.docx': extract_docx,
    '.doc': extract_doc,
    '.odt': extract_odt,
    '.ods': extract_ods,
    '.odp': extract_odp,
    '.rtf': extract_rtf,
    '.pptx': extract_pptx,
    '.xlsx': extract_xlsx,
    '.xls': extract_xls,
    '.csv': extract_csv,
    '.html': extract_html,
    '.htm': extract_html,
}

IMAGE_EXTS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.tif', '.webp'}

def detect_by_mime(path: Path) -> Optional[str]:
    """Fallback: detecta por MIME type se extensao desconhecida."""
    try:
        import magic
        mime = magic.from_file(str(path), mime=True)
        mapping = {
            'application/pdf': '.pdf',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
            'application/msword': '.doc',
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': '.xlsx',
            'application/vnd.ms-excel': '.xls',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
            'application/vnd.oasis.opendocument.text': '.odt',
            'application/vnd.oasis.opendocument.spreadsheet': '.ods',
            'application/vnd.oasis.opendocument.presentation': '.odp',
            'text/rtf': '.rtf',
            'text/csv': '.csv',
            'text/html': '.html',
            'image/png': '.png',
            'image/jpeg': '.jpg',
            'image/gif': '.gif',
            'image/bmp': '.bmp',
            'image/tiff': '.tiff',
            'image/webp': '.webp',
            'text/plain': '.txt',
        }
        return mapping.get(mime)
    except ImportError:
        return None
    except Exception:
        return None

# ============================================================
# Leitura principal
# ============================================================

def read_file(filepath: str, max_chars: int = 100000) -> dict:
    path = Path(filepath)
    if not path.exists():
        return {'error': f'Arquivo nao encontrado: {filepath}', 'text': ''}

    ext = path.suffix.lower()
    detected_ext = ext

    try:
        # Extractor direto por extensao
        if ext in EXTRACTORS:
            text = EXTRACTORS[ext](path)
        elif ext in IMAGE_EXTS:
            text = extract_image(path)
        elif ext in TEXT_EXTS:
            text = path.read_text(encoding='utf-8', errors='replace')
        else:
            # Fallback: detecta por MIME
            mime_ext = detect_by_mime(path)
            if mime_ext and mime_ext in EXTRACTORS:
                detected_ext = mime_ext
                text = EXTRACTORS[mime_ext](path)
            elif mime_ext in IMAGE_EXTS:
                detected_ext = mime_ext
                text = extract_image(path)
            else:
                # Ultimo recurso: tenta ler como texto
                try:
                    text = path.read_text(encoding='utf-8', errors='replace')
                    # Valida se tem conteudo legivel
                    if len(text.strip()) > 0 and sum(1 for c in text if c.isprintable()) > len(text) * 0.5:
                        pass
                    else:
                        text = f'[Formato nao suportado: {ext or mime_ext or "desconhecido"}]'
                except Exception:
                    text = f'[Formato nao suportado: {ext or "desconhecido"}]'
    except ImportError as e:
        missing = str(e).split("'")[1] if "'" in str(e) else str(e)
        return {
            'error': f'Dependencia faltando: {missing}. Rode: pip install -r scripts/requirements.txt',
            'text': '', 'file': filepath,
        }
    except Exception as e:
        return {'error': str(e), 'text': '', 'file': filepath}

    truncated = len(text) > max_chars
    return {
        'file': filepath,
        'name': path.name,
        'extension': detected_ext,
        'size_chars': len(text),
        'size_bytes': path.stat().st_size if path.exists() else 0,
        'truncated': truncated,
        'text': text[:max_chars],
    }

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print('Uso: python3 scripts/file-reader.py <arquivo> [max_chars]')
        sys.exit(1)
    max_chars = int(sys.argv[2]) if len(sys.argv) > 2 else 100000
    result = read_file(sys.argv[1], max_chars)
    print(json.dumps(result, ensure_ascii=False, indent=2))
